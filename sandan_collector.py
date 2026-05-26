import argparse
import csv
import hashlib
import json
import logging
import os
import re
import subprocess
import sys
import time
import zipfile
import zlib
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional
from urllib.parse import unquote, urlencode, urljoin

import requests
from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from tqdm import tqdm

try:
    import pikepdf
except ImportError:
    pikepdf = None

try:
    import olefile
except ImportError:
    olefile = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from sandan_rag.collection_state import SQLiteCollectionState
except Exception:
    SQLiteCollectionState = None


# Windows PowerShell / CMD may use CP949/GBK by default.
# Reconfigure stdout/stderr so Korean log messages do not crash collection.
def configure_utf8_stdio() -> None:
    for stream_name in ["stdout", "stderr"]:
        stream = getattr(sys, stream_name, None)
        try:
            if hasattr(stream, "reconfigure"):
                stream.reconfigure(encoding="utf-8", errors="replace")
        except Exception:
            pass


configure_utf8_stdio()

BASE_URL = "https://research.khu.ac.kr"
BOARD_CODE = "BMSR00040"
LIST_ENDPOINT = f"{BASE_URL}/research/user/bbs/{BOARD_CODE}/list.do"

TARGET_MENUS = [
    {
        "menu_no": "6400049",
        "menu_name": "대외연구비_규정지침",
        "list_url": "https://research.khu.ac.kr/research/user/bbs/BMSR00040/list.do?menuNo=6400049",
    },
    {
        "menu_no": "6400050",
        "menu_name": "자료실_서식양식",
        "list_url": "https://research.khu.ac.kr/research/user/bbs/BMSR00040/list.do?menuNo=6400050",
    },
    {
        "menu_no": "6400163",
        "menu_name": "산학협력단_규정지침",
        "list_url": "https://research.khu.ac.kr/research/user/bbs/BMSR00040/list.do?menuNo=6400163",
    },
]

USER_AGENT = (
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
    "AppleWebKit/537.36 (KHTML, like Gecko) "
    "Chrome/126.0.0.0 Safari/537.36"
)


def clean_text(text: str) -> str:
    text = text or ""
    text = text.replace("\xa0", " ")
    text = text.replace("\ufeff", " ")
    text = re.sub(r"\r\n|\r", "\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def safe_filename(text: str, max_len: int = 150) -> str:
    text = clean_text(text)
    text = re.sub(r"[\\/:*?\"<>|]", "_", text)
    text = re.sub(r"\s+", " ", text).strip()

    if not text:
        return "untitled"

    return text[:max_len]


def sha256_short(data: bytes | str) -> str:
    if isinstance(data, str):
        data = data.encode("utf-8", errors="ignore")

    return hashlib.sha256(data).hexdigest()[:16]


def decode_header_filename(raw: str) -> str:
    if not raw:
        return ""

    raw = raw.strip().strip("\"")

    for encoding in ["utf-8", "cp949", "euc-kr"]:
        try:
            return unquote(raw, encoding=encoding)
        except Exception:
            pass

    for encoding in ["utf-8", "cp949", "euc-kr"]:
        try:
            return raw.encode("latin1").decode(encoding)
        except Exception:
            pass

    return raw


def filename_from_content_disposition(value: str) -> str:
    if not value:
        return ""

    matched = re.search(r"filename\*=UTF-8''([^;]+)", value, flags=re.I)
    if matched:
        return decode_header_filename(matched.group(1))

    matched = re.search(r"filename\*=[^']*''([^;]+)", value, flags=re.I)
    if matched:
        return decode_header_filename(matched.group(1))

    matched = re.search(r"filename=\"?([^\";]+)\"?", value, flags=re.I)
    if matched:
        return decode_header_filename(matched.group(1))

    return ""


def guess_extension(file_name: str, content_type: str = "") -> str:
    suffix = Path(file_name).suffix.lower()

    if suffix:
        return suffix

    content_type = (content_type or "").lower()

    if "pdf" in content_type:
        return ".pdf"
    if "hwp" in content_type:
        return ".hwp"
    if "word" in content_type:
        return ".docx"
    if "powerpoint" in content_type or "presentation" in content_type:
        return ".pptx"
    if "excel" in content_type or "spreadsheet" in content_type:
        return ".xlsx"
    if "zip" in content_type:
        return ".zip"
    if "text" in content_type:
        return ".txt"

    return ".bin"


class SandanAttachmentCollector:
    def __init__(
        self,
        output_dir: str = "data/sandan_attachment_kb",
        sleep_sec: float = 0.35,
        timeout: int = 40,
        max_pages: int = 300,
        full: bool = False,
        skip_existing_files: bool = True,
        suppress_pdf_warnings: bool = True,
        repair_pdf: bool = True,
        state_db_path: str | None = None,
    ):
        self.output_dir = Path(output_dir)
        self.html_dir = self.output_dir / "raw_html"
        self.attach_dir = self.output_dir / "attachments"
        self.text_dir = self.output_dir / "attachment_texts"
        self.failed_dir = self.output_dir / "failed"
        self.zip_extract_dir = self.output_dir / "zip_extracted"
        self.hwp5txt_dir = self.output_dir / "hwp5txt_output"
        self.ppt_convert_dir = self.output_dir / "ppt_converted"
        self.pdf_fixed_dir = self.output_dir / "pdf_fixed"

        self.html_dir.mkdir(parents=True, exist_ok=True)
        self.attach_dir.mkdir(parents=True, exist_ok=True)
        self.text_dir.mkdir(parents=True, exist_ok=True)
        self.failed_dir.mkdir(parents=True, exist_ok=True)
        self.zip_extract_dir.mkdir(parents=True, exist_ok=True)
        self.hwp5txt_dir.mkdir(parents=True, exist_ok=True)
        self.ppt_convert_dir.mkdir(parents=True, exist_ok=True)
        self.pdf_fixed_dir.mkdir(parents=True, exist_ok=True)

        self.sleep_sec = sleep_sec
        self.timeout = timeout
        self.max_pages = max_pages
        self.full = full
        self.skip_existing_files = skip_existing_files
        self.suppress_pdf_warnings = suppress_pdf_warnings
        self.repair_pdf = repair_pdf

        if self.suppress_pdf_warnings:
            logging.getLogger("pypdf").setLevel(logging.ERROR)
            logging.getLogger("pypdf._reader").setLevel(logging.ERROR)
        else:
            logging.getLogger("pypdf").setLevel(logging.WARNING)
            logging.getLogger("pypdf._reader").setLevel(logging.WARNING)

        self.records_jsonl = self.output_dir / "sandan_attachment_records.jsonl"
        self.records_csv = self.output_dir / "sandan_attachment_records.csv"
        self.merged_txt = self.output_dir / "sandan_attachment_merged_text.txt"
        self.failed_jsonl = self.failed_dir / "failed_downloads.jsonl"

        state_db_value = state_db_path or os.getenv("SANDAN_STATE_DB_PATH", "data/sandan_collection_state.sqlite3")
        self.state_db_path = Path(state_db_value)
        self.collection_state = None
        if SQLiteCollectionState is not None:
            try:
                self.collection_state = SQLiteCollectionState(self.state_db_path)
            except Exception as exc:
                print(f"[WARN] collection state DB disabled: {self.state_db_path}, error={exc}")

        self.session = requests.Session()
        self.session.headers.update(
            {
                "User-Agent": USER_AGENT,
                "Accept-Language": "ko-KR,ko;q=0.9,en-US;q=0.8,en;q=0.7,zh-CN;q=0.6",
                "Referer": LIST_ENDPOINT,
            }
        )

    def get(self, url: str, referer: Optional[str] = None) -> requests.Response:
        headers = {}

        if referer:
            headers["Referer"] = referer

        response = self.session.get(
            url,
            timeout=self.timeout,
            allow_redirects=True,
            headers=headers,
        )
        response.raise_for_status()

        if not response.encoding or response.encoding.lower() == "iso-8859-1":
            response.encoding = response.apparent_encoding or "utf-8"

        time.sleep(self.sleep_sec)
        return response

    def build_list_url(self, menu_no: str, page_index: int) -> str:
        params = {
            "menuNo": menu_no,
            "boardType": "",
            "pageIndex": str(page_index),
            "searchCondition": "",
            "searchKeyword": "",
            "userDisplayBbsType": "D",
        }

        return f"{LIST_ENDPOINT}?{urlencode(params)}"

    def build_detail_url(self, menu_no: str, board_id: str) -> str:
        params = {
            "menuNo": menu_no,
            "boardId": board_id,
        }

        return f"{BASE_URL}/research/user/bbs/{BOARD_CODE}/view.do?{urlencode(params)}"

    def parse_last_page(self, html: str) -> int:
        page_numbers = []

        for matched in re.findall(r"fnSubmitForm\((\d+)\)", html):
            try:
                page_numbers.append(int(matched))
            except ValueError:
                pass

        if page_numbers:
            return max(page_numbers)

        total_match = re.search(r"전체\s*<strong>\s*([0-9,]+)\s*</strong>", html)
        if total_match:
            total_count = int(total_match.group(1).replace(",", ""))
            return max(1, (total_count + 9) // 10)

        total_match = re.search(r"전체\s*([0-9,]+)\s*건", clean_text(html))
        if total_match:
            total_count = int(total_match.group(1).replace(",", ""))
            return max(1, (total_count + 9) // 10)

        return 1

    def parse_list_page(self, html: str, menu: Dict) -> List[Dict]:
        soup = BeautifulSoup(html, "html.parser")
        results = []

        for row in soup.select("table tr"):
            row_html = str(row)
            board_match = re.search(r"javascript:view\(['\"]?(\d+)['\"]?\)", row_html)

            if not board_match:
                continue

            board_id = board_match.group(1)

            title_tag = row.select_one("td.tal a")
            if title_tag is None:
                title_tag = row.find("a")

            post_title = clean_text(title_tag.get_text(" ", strip=True)) if title_tag else ""
            post_title = re.sub(r"\s+", " ", post_title).strip()

            registered_date = ""
            cells = row.find_all("td")

            for cell in reversed(cells):
                cell_text = clean_text(cell.get_text(" ", strip=True))
                date_match = re.search(r"\d{4}-\d{2}-\d{2}", cell_text)
                if date_match:
                    registered_date = date_match.group(0)
                    break

            file_cell = row.select_one("td.file")
            has_attachment = False

            if file_cell is not None:
                has_attachment = bool(file_cell.select("img[alt*='첨부파일'], img[alt*='첨부']"))

            results.append(
                {
                    "menu_no": menu["menu_no"],
                    "menu_name": menu["menu_name"],
                    "board_id": board_id,
                    "post_uid": f"{menu['menu_no']}_{board_id}",
                    "post_title": post_title,
                    "registered_date": registered_date,
                    "detail_url": self.build_detail_url(menu["menu_no"], board_id),
                    "has_attachment_icon": has_attachment,
                    "row_text": clean_text(row.get_text(" ", strip=True)),
                }
            )

        unique = {}
        for item in results:
            unique[item["post_uid"]] = item

        return list(unique.values())

    def collect_post_list_for_menu(self, menu: Dict) -> List[Dict]:
        menu_no = menu["menu_no"]
        menu_name = menu["menu_name"]

        first_url = self.build_list_url(menu_no, 1)
        first_response = self.get(first_url)

        menu_html_dir = self.html_dir / safe_filename(f"{menu_no}_{menu_name}", 120)
        menu_html_dir.mkdir(parents=True, exist_ok=True)

        first_html_path = menu_html_dir / "list_page_001.html"
        first_html_path.write_text(first_response.text, encoding="utf-8")

        last_page = self.parse_last_page(first_response.text)
        last_page = min(last_page, self.max_pages)

        print(f"[MENU] {menu_no} {menu_name}")
        print(f"[INFO] detected last page: {last_page}")

        all_posts = {}

        for page_index in range(1, last_page + 1):
            if page_index == 1:
                html = first_response.text
            else:
                url = self.build_list_url(menu_no, page_index)
                response = self.get(url)
                html = response.text

                html_path = menu_html_dir / f"list_page_{page_index:03d}.html"
                html_path.write_text(html, encoding="utf-8")

            posts = self.parse_list_page(html, menu)
            new_count = 0

            for post in posts:
                post_uid = post["post_uid"]

                if post_uid not in all_posts:
                    all_posts[post_uid] = post
                    new_count += 1

            print(f"[LIST] menu={menu_no}, page={page_index}, found={len(posts)}, new={new_count}")

        posts = list(all_posts.values())
        posts.sort(
            key=lambda item: (
                item.get("registered_date", ""),
                item.get("menu_no", ""),
                item.get("board_id", ""),
            ),
            reverse=True,
        )

        return posts

    def collect_all_posts(self) -> List[Dict]:
        all_posts = {}

        for menu in TARGET_MENUS:
            posts = self.collect_post_list_for_menu(menu)

            for post in posts:
                all_posts[post["post_uid"]] = post

        posts = list(all_posts.values())
        posts.sort(
            key=lambda item: (
                item.get("registered_date", ""),
                item.get("menu_no", ""),
                item.get("board_id", ""),
            ),
            reverse=True,
        )

        return posts

    def load_records_from_jsonl(self) -> List[Dict]:
        if not self.records_jsonl.exists():
            return []

        records = []

        with self.records_jsonl.open("r", encoding="utf-8") as f:
            for line in f:
                if not line.strip():
                    continue
                try:
                    item = json.loads(line)
                    if isinstance(item, dict) and item.get("attachment_key"):
                        records.append(item)
                except Exception:
                    continue

        return records

    def load_records_from_state_db(self) -> List[Dict]:
        if self.collection_state is None:
            return []
        try:
            return self.collection_state.load_records()
        except Exception as exc:
            print(f"[WARN] failed to read collection state DB: {exc}")
            return []

    def load_existing_records(self) -> List[Dict]:
        if self.full:
            return []

        records = self.load_records_from_state_db() + self.load_records_from_jsonl()
        dedup = {}
        for record in records:
            key = record.get("attachment_key", "")
            if key:
                dedup[key] = record
        return list(dedup.values())

    def load_existing_attachment_keys(self) -> set:
        if self.full:
            return set()

        keys = set()

        if self.collection_state is not None:
            try:
                keys.update(self.collection_state.keys())
            except Exception as exc:
                print(f"[WARN] failed to read keys from collection state DB: {exc}")

        for item in self.load_records_from_jsonl():
            key = item.get("attachment_key", "")
            if key:
                keys.add(key)

        return keys

    def persist_record_incremental(self, record: Dict) -> None:
        """Save one processed attachment immediately.

        This prevents losing all progress if the Streamlit process or WebSocket dies
        before the final save_records() step.
        """
        if self.collection_state is not None:
            try:
                self.collection_state.upsert_record(record)
            except Exception as exc:
                print(f"[WARN] failed to upsert collection state: {exc}")

        try:
            with self.records_jsonl.open("a", encoding="utf-8") as f:
                f.write(json.dumps(record, ensure_ascii=False) + "\n")
        except Exception as exc:
            print(f"[WARN] failed to append record jsonl: {exc}")

    def extract_detail_metadata(self, soup: BeautifulSoup, fallback: Dict) -> Dict:
        article = soup.select_one("article.bbs-view")

        title = fallback.get("post_title", "")
        registered_date = fallback.get("registered_date", "")
        author = ""

        if article:
            title_tag = article.select_one("header.top h2.t")
            if title_tag:
                title = clean_text(title_tag.get_text(" ", strip=True))
                title = title.replace("[]", "")
                title = title.replace("&nbsp;", " ")
                title = re.sub(r"\s+", " ", title).strip()

            full_text = clean_text(article.get_text("\n", strip=True))
        else:
            full_text = clean_text(soup.get_text("\n", strip=True))

        date_match = re.search(
            r"등록일\s*([0-9]{4}-[0-9]{2}-[0-9]{2}(?:\s+[0-9:.]+)?)",
            full_text,
        )
        if date_match:
            registered_date = date_match.group(1).strip()

        author_match = re.search(r"작성자\s*([^\n]+)", full_text)
        if author_match:
            author = clean_text(author_match.group(1))

        return {
            "post_title": title,
            "registered_date": registered_date,
            "author": author,
        }

    def extract_post_body_text(self, soup: BeautifulSoup) -> str:
        article = soup.select_one("article.bbs-view")
        if not article:
            return ""

        body = article.select_one(".bbs-view_c")
        if not body:
            return ""

        body_copy = BeautifulSoup(str(body), "html.parser")

        for tag in body_copy.select("script, style, noscript"):
            tag.decompose()

        return clean_text(body_copy.get_text("\n", strip=True))

    def extract_attachment_links(self, soup: BeautifulSoup, detail_url: str) -> List[Dict]:
        article = soup.select_one("article.bbs-view")

        if article is None:
            return []

        attachments = []

        for a in article.select("a[href*='fileDown.do']"):
            href = a.get("href", "") or ""
            name = clean_text(a.get_text(" ", strip=True))

            if not href:
                continue

            lower_href = href.lower()

            if "imagesrc.do" in lower_href:
                continue

            if "filedown.do" not in lower_href:
                continue

            file_url = urljoin(BASE_URL, href.replace("&amp;", "&"))

            attachments.append(
                {
                    "name": name or "attachment",
                    "url": file_url,
                    "source": "article_bbs_view",
                }
            )

        unique = {}

        for item in attachments:
            unique[item["url"]] = item

        return list(unique.values())

    def validate_download_response(self, response: requests.Response) -> bool:
        if response.status_code != 200:
            return False

        if not response.content:
            return False

        content_type = response.headers.get("Content-Type", "").lower()
        content_disposition = response.headers.get("Content-Disposition", "").lower()

        if "text/html" in content_type and "attachment" not in content_disposition:
            preview = response.text[:1000].lower()
            if "<html" in preview or "error" in preview or "페이지" in preview:
                return False

        return True

    def save_failed_download(self, post: Dict, attachment: Dict, reason: str) -> None:
        payload = {
            "menu_no": post.get("menu_no", ""),
            "menu_name": post.get("menu_name", ""),
            "board_id": post.get("board_id", ""),
            "post_title": post.get("post_title", ""),
            "registered_date": post.get("registered_date", ""),
            "detail_url": post.get("detail_url", ""),
            "attachment_name": attachment.get("name", ""),
            "attachment_url": attachment.get("url", ""),
            "reason": reason,
            "failed_at": datetime.now().isoformat(timespec="seconds"),
        }

        with self.failed_jsonl.open("a", encoding="utf-8") as f:
            f.write(json.dumps(payload, ensure_ascii=False) + "\n")

    def download_attachment(self, attachment: Dict, post: Dict, index: int) -> Optional[Dict]:
        url = attachment["url"]
        menu_no = post["menu_no"]
        menu_name = post["menu_name"]
        board_id = post["board_id"]
        detail_url = post.get("detail_url", "")

        try:
            response = self.get(url, referer=detail_url or LIST_ENDPOINT)
        except Exception as exc:
            self.save_failed_download(post, attachment, f"request failed: {exc}")
            print(f"[WARN] attachment request failed: menu={menu_no}, board_id={board_id}, url={url}, error={exc}")
            return None

        if not self.validate_download_response(response):
            self.save_failed_download(post, attachment, "invalid download response")
            print(f"[WARN] invalid download response: menu={menu_no}, board_id={board_id}, url={url}")
            return None

        content_type = response.headers.get("Content-Type", "")
        content_disposition = response.headers.get("Content-Disposition", "")
        server_name = filename_from_content_disposition(content_disposition)

        name_hint = server_name or attachment.get("name", "") or f"attachment_{index}"
        name_hint = safe_filename(name_hint, 120)
        extension = guess_extension(name_hint, content_type)

        if not Path(name_hint).suffix:
            name_hint = f"{name_hint}{extension}"

        menu_prefix = safe_filename(f"{menu_no}_{menu_name}", 80)
        file_name = safe_filename(f"{menu_prefix}_{board_id}_{index:02d}_{name_hint}", 220)

        if not Path(file_name).suffix:
            file_name = f"{file_name}{extension}"

        file_path = self.attach_dir / menu_prefix / file_name
        file_path.parent.mkdir(parents=True, exist_ok=True)

        new_hash = sha256_short(response.content)

        if file_path.exists() and self.skip_existing_files:
            old_hash = sha256_short(file_path.read_bytes())

            if old_hash == new_hash:
                return {
                    "attachment_name": file_path.name,
                    "attachment_path": str(file_path),
                    "attachment_url": response.url,
                    "content_type": content_type,
                    "file_hash": old_hash,
                }

        file_path.write_bytes(response.content)

        return {
            "attachment_name": file_path.name,
            "attachment_path": str(file_path),
            "attachment_url": response.url,
            "content_type": content_type,
            "file_hash": new_hash,
        }

    def is_garbled_text(self, text: str) -> bool:
        text = clean_text(text)

        if not text:
            return True

        sample = text[:5000]

        korean_count = len(re.findall(r"[가-힣]", sample))
        english_number_count = len(re.findall(r"[A-Za-z0-9]", sample))
        normal_symbol_count = len(re.findall(r"[\s.,:;()\[\]{}<>/%\-_=+·ㆍ※①-⑳Ⅰ-Ⅹ]", sample))

        normal_count = korean_count + english_number_count + normal_symbol_count
        total_count = max(len(sample), 1)
        normal_ratio = normal_count / total_count

        strange_count = len(re.findall(r"[\u0370-\u03ff\u0600-\u06ff\u0800-\u0fff\u0c00-\u0cff]", sample))
        replacement_count = sample.count("�")

        if korean_count < 20 and normal_ratio < 0.5:
            return True

        if strange_count > 80:
            return True

        if replacement_count > 30:
            return True

        if normal_ratio < 0.35:
            return True

        return False

    def read_text_file_with_encoding_guess(self, path: Path) -> str:
        encodings = ["utf-8", "utf-16", "utf-16le", "cp949", "euc-kr"]

        for encoding in encodings:
            try:
                text = path.read_text(encoding=encoding, errors="ignore")
                text = clean_text(text)

                if text and not self.is_garbled_text(text):
                    return text

            except Exception:
                continue

        return ""

    def extract_hwp_with_hwp5txt(self, path: Path) -> str:
        out_name = safe_filename(f"{path.stem}_{sha256_short(str(path))}", 160)
        out_path = self.hwp5txt_dir / f"{out_name}.txt"

        if out_path.exists():
            try:
                out_path.unlink()
            except Exception:
                pass

        command = ["hwp5txt", "--output", str(out_path), str(path)]

        try:
            result = subprocess.run(
                command,
                capture_output=True,
                text=True,
                encoding="utf-8",
                errors="ignore",
                timeout=180,
                check=False,
            )

            if out_path.exists():
                text = self.read_text_file_with_encoding_guess(out_path)

                if text and not self.is_garbled_text(text):
                    return text

            stdout_text = clean_text(result.stdout)

            if stdout_text and not self.is_garbled_text(stdout_text):
                return stdout_text

            if result.stderr.strip():
                print(f"[WARN] hwp5txt stderr: {path.name}: {result.stderr.strip()[:300]}")

        except Exception as exc:
            print(f"[WARN] hwp5txt failed: {path}, error={exc}")

        return ""

    def extract_hwp_raw_fallback(self, path: Path) -> str:
        if olefile is None:
            return ""

        if not olefile.isOleFile(str(path)):
            return ""

        try:
            hwp = olefile.OleFileIO(str(path))
            streams = hwp.listdir()

            section_streams = [
                stream for stream in streams
                if len(stream) >= 2
                and stream[0] == "BodyText"
                and stream[1].startswith("Section")
            ]

            section_streams = sorted(section_streams, key=lambda x: x[1])
            compressed = self.is_hwp_compressed(hwp)

            parts = []

            for section in section_streams:
                data = hwp.openstream(section).read()

                if compressed:
                    try:
                        data = zlib.decompress(data, -15)
                    except Exception:
                        pass

                decoded = data.decode("utf-16le", errors="ignore")

                tokens = re.findall(
                    r"[가-힣A-Za-z0-9][가-힣A-Za-z0-9\s.,:;()\[\]{}<>/%\-_=+·ㆍ※①-⑳Ⅰ-Ⅹ]{2,}",
                    decoded,
                )

                text = clean_text("\n".join(tokens))

                if text:
                    parts.append(text)

            output = clean_text("\n\n".join(parts))

            if self.is_garbled_text(output):
                return ""

            return output

        except Exception:
            return ""

    def extract_hwp(self, path: Path) -> str:
        text = self.extract_hwp_with_hwp5txt(path)

        if text:
            return text

        text = self.extract_hwp_raw_fallback(path)

        if text:
            return text

        print(f"[WARN] HWP text extraction failed or garbled: {path}")
        return ""

    def is_hwp_compressed(self, hwp) -> bool:
        try:
            header = hwp.openstream("FileHeader").read()
            flags = int.from_bytes(header[36:40], "little")
            return bool(flags & 1)
        except Exception:
            return True

    def extract_text_from_file(self, file_path: str, zip_depth: int = 0) -> str:
        path = Path(file_path)
        suffix = path.suffix.lower()

        try:
            if suffix == ".pdf":
                return self.extract_pdf(path)
            if suffix == ".docx":
                return self.extract_docx(path)
            if suffix in [".pptx", ".ppt"]:
                return self.extract_ppt(path)
            if suffix in [".xlsx", ".xlsm"]:
                return self.extract_xlsx(path)
            if suffix == ".hwpx":
                return self.extract_hwpx(path)
            if suffix == ".hwp":
                return self.extract_hwp(path)
            if suffix in [".txt", ".md", ".csv"]:
                return self.read_text_file_with_encoding_guess(path) or path.read_text(encoding="utf-8", errors="ignore")
            if suffix == ".zip":
                return self.extract_zip(path, zip_depth=zip_depth)
        except Exception as exc:
            print(f"[WARN] text extraction failed: {file_path}, error={exc}")

        return ""

    def is_probably_pdf(self, path: Path) -> bool:
        try:
            if not path.exists() or path.stat().st_size == 0:
                return False

            with path.open("rb") as f:
                header = f.read(5)

            return header == b"%PDF-"

        except Exception:
            return False

    def repair_pdf_with_pikepdf(self, path: Path) -> Optional[Path]:
        if not self.repair_pdf:
            return None

        if pikepdf is None:
            return None

        fixed_name = safe_filename(f"{path.stem}_{sha256_short(str(path))}_fixed", 180)
        fixed_path = self.pdf_fixed_dir / f"{fixed_name}.pdf"

        try:
            with pikepdf.open(str(path)) as pdf:
                pdf.save(str(fixed_path))

            if fixed_path.exists() and fixed_path.stat().st_size > 0:
                print(f"[PDF] repaired: {path.name} -> {fixed_path.name}")
                return fixed_path

        except Exception as exc:
            print(f"[WARN] PDF repair failed: {path}, error={exc}")

        return None

    def _extract_pdf_once(self, path: Path) -> str:
        reader = PdfReader(str(path), strict=False)
        pages = []

        for i, page in enumerate(reader.pages):
            try:
                text = page.extract_text() or ""
                text = clean_text(text)

                if text:
                    pages.append(f"[PDF Page {i + 1}]\n{text}")

            except Exception as exc:
                print(f"[WARN] PDF page extraction failed: {path}, page={i + 1}, error={exc}")
                continue

        return clean_text("\n\n".join(pages))

    def extract_pdf(self, path: Path) -> str:
        try:
            if not path.exists() or path.stat().st_size == 0:
                print(f"[WARN] empty PDF file: {path}")
                return ""

            if not self.is_probably_pdf(path):
                print(f"[WARN] not a real PDF file or invalid PDF header: {path}")
                return ""

            try:
                text = self._extract_pdf_once(path)

                if text:
                    return text

                fixed_path = self.repair_pdf_with_pikepdf(path)

                if fixed_path is not None:
                    return self._extract_pdf_once(fixed_path)

                return ""

            except Exception as exc:
                print(f"[WARN] PDF first extraction failed, trying repair: {path}, error={exc}")

                fixed_path = self.repair_pdf_with_pikepdf(path)

                if fixed_path is not None:
                    try:
                        return self._extract_pdf_once(fixed_path)
                    except Exception as second_exc:
                        print(f"[WARN] PDF extraction failed after repair: {path}, error={second_exc}")

                return ""

        except Exception as exc:
            print(f"[WARN] PDF extraction failed: {path}, error={exc}")
            return ""

    def extract_docx(self, path: Path) -> str:
        doc = Document(str(path))
        parts = []

        for paragraph in doc.paragraphs:
            text = clean_text(paragraph.text)
            if text:
                parts.append(text)

        for table in doc.tables:
            for row in table.rows:
                values = [clean_text(cell.text) for cell in row.cells]
                values = [value for value in values if value]
                if values:
                    parts.append(" | ".join(values))

        return clean_text("\n".join(parts))

    def extract_pptx(self, path: Path) -> str:
        if Presentation is None:
            print("[WARN] python-pptx not installed. Cannot extract PPTX.")
            return ""

        prs = Presentation(str(path))
        parts = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_parts = []

            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = clean_text(shape.text)
                    if text:
                        slide_parts.append(text)

                if getattr(shape, "has_table", False):
                    table = shape.table

                    for row in table.rows:
                        values = []

                        for cell in row.cells:
                            cell_text = clean_text(cell.text)
                            if cell_text:
                                values.append(cell_text)

                        if values:
                            slide_parts.append(" | ".join(values))

            if getattr(slide, "has_notes_slide", False):
                try:
                    notes_text = clean_text(slide.notes_slide.notes_text_frame.text)
                    if notes_text:
                        slide_parts.append(f"[Notes]\n{notes_text}")
                except Exception:
                    pass

            if slide_parts:
                parts.append(
                    "\n".join(
                        [
                            f"[PPT Slide {slide_idx}]",
                            "\n".join(slide_parts),
                        ]
                    )
                )

        return clean_text("\n\n".join(parts))

    def find_soffice_path(self) -> str:
        candidates = [
            "soffice",
            "libreoffice",
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]

        for candidate in candidates:
            try:
                result = subprocess.run(
                    [candidate, "--version"],
                    capture_output=True,
                    text=True,
                    timeout=10,
                    check=False,
                )

                if result.returncode == 0:
                    return candidate

            except Exception:
                continue

        return ""

    def extract_ppt_with_libreoffice(self, path: Path) -> str:
        soffice = self.find_soffice_path()

        if not soffice:
            print(f"[WARN] LibreOffice not found. Cannot convert PPT: {path}")
            return ""

        out_dir = self.ppt_convert_dir / safe_filename(f"{path.stem}_{sha256_short(str(path))}", 180)
        out_dir.mkdir(parents=True, exist_ok=True)

        try:
            subprocess.run(
                [
                    soffice,
                    "--headless",
                    "--convert-to",
                    "pptx",
                    "--outdir",
                    str(out_dir),
                    str(path),
                ],
                capture_output=True,
                text=True,
                timeout=180,
                check=False,
            )

            pptx_candidates = list(out_dir.glob("*.pptx"))

            for pptx_path in pptx_candidates:
                text = self.extract_pptx(pptx_path)
                text = clean_text(text)

                if text:
                    return text

        except Exception as exc:
            print(f"[WARN] PPT conversion failed: {path}, error={exc}")

        return ""

    def extract_ppt(self, path: Path) -> str:
        suffix = path.suffix.lower()

        if suffix == ".pptx":
            return self.extract_pptx(path)

        if suffix == ".ppt":
            return self.extract_ppt_with_libreoffice(path)

        return ""

    def extract_xlsx(self, path: Path) -> str:
        workbook = load_workbook(str(path), read_only=True, data_only=True)
        parts = []

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            parts.append(f"[Sheet] {sheet_name}")

            for row in sheet.iter_rows(values_only=True):
                values = []

                for value in row:
                    if value is None:
                        continue

                    text = clean_text(str(value))
                    if text:
                        values.append(text)

                if values:
                    parts.append(" | ".join(values))

        return clean_text("\n".join(parts))

    def extract_hwpx(self, path: Path) -> str:
        parts = []

        with zipfile.ZipFile(path, "r") as zf:
            xml_names = [
                name for name in zf.namelist()
                if name.lower().endswith(".xml")
                and (
                    "contents/" in name.lower()
                    or "section" in name.lower()
                    or "bodytext" in name.lower()
                )
            ]

            for name in xml_names:
                try:
                    raw = zf.read(name).decode("utf-8", errors="ignore")
                    raw = re.sub(r"<[^>]+>", " ", raw)
                    raw = clean_text(raw)

                    if raw:
                        parts.append(raw)
                except Exception:
                    continue

        return clean_text("\n\n".join(parts))

    def decode_zip_member_name(self, info: zipfile.ZipInfo) -> str:
        name = info.filename

        if info.flag_bits & 0x800:
            return name

        try:
            return name.encode("cp437").decode("cp949")
        except Exception:
            pass

        try:
            return name.encode("cp437").decode("euc-kr")
        except Exception:
            pass

        return name

    def extract_zip(self, path: Path, zip_depth: int = 0) -> str:
        if zip_depth > 3:
            return ""

        parts = []

        zip_hash = sha256_short(str(path))
        extract_root = self.zip_extract_dir / safe_filename(f"{path.stem}_{zip_hash}", 180)
        extract_root.mkdir(parents=True, exist_ok=True)

        supported_suffixes = {
            ".pdf",
            ".docx",
            ".pptx",
            ".ppt",
            ".xlsx",
            ".xlsm",
            ".hwp",
            ".hwpx",
            ".txt",
            ".md",
            ".csv",
            ".zip",
        }

        try:
            with zipfile.ZipFile(path, "r") as zf:
                for info in zf.infolist():
                    if info.is_dir():
                        continue

                    decoded_name = self.decode_zip_member_name(info)
                    suffix = Path(decoded_name).suffix.lower()

                    if suffix not in supported_suffixes:
                        continue

                    safe_member_name = safe_filename(decoded_name.replace("\\", "_").replace("/", "_"), 220)
                    extracted_path = extract_root / safe_member_name

                    counter = 1
                    while extracted_path.exists():
                        extracted_path = extract_root / f"{extracted_path.stem}_{counter}{extracted_path.suffix}"
                        counter += 1

                    with zf.open(info) as source:
                        extracted_path.write_bytes(source.read())

                    text = self.extract_text_from_file(str(extracted_path), zip_depth=zip_depth + 1)
                    text = clean_text(text)

                    if text and self.is_garbled_text(text):
                        print(f"[WARN] garbled text skipped inside zip: {decoded_name}")
                        text = ""

                    if text:
                        parts.append(
                            "\n".join(
                                [
                                    f"[ZIP 내부파일 경로]\n{decoded_name}",
                                    "[ZIP 내부파일 본문]",
                                    text,
                                ]
                            )
                        )
                    else:
                        parts.append(
                            "\n".join(
                                [
                                    f"[ZIP 내부파일 경로]\n{decoded_name}",
                                    "[ZIP 내부파일 본문]",
                                    "",
                                ]
                            )
                        )

        except Exception as exc:
            print(f"[WARN] ZIP extraction failed: {path}, error={exc}")
            return ""

        return clean_text("\n\n".join(parts))

    def collect_detail_attachments(self, post: Dict, existing_keys: set) -> List[Dict]:
        detail_url = post["detail_url"]
        menu_no = post["menu_no"]
        menu_name = post["menu_name"]
        board_id = post["board_id"]
        post_uid = post["post_uid"]

        try:
            response = self.get(detail_url)
        except Exception as exc:
            print(f"[WARN] detail failed: menu={menu_no}, board_id={board_id}, error={exc}")
            return []

        detail_html_dir = self.html_dir / safe_filename(f"{menu_no}_{menu_name}", 120)
        detail_html_dir.mkdir(parents=True, exist_ok=True)

        detail_html_path = detail_html_dir / f"detail_{board_id}.html"
        detail_html_path.write_text(response.text, encoding="utf-8")

        soup = BeautifulSoup(response.text, "html.parser")
        detail_meta = self.extract_detail_metadata(soup, post)
        post_body_text = self.extract_post_body_text(soup)
        attachments = self.extract_attachment_links(soup, detail_url)

        print(
            f"[DETAIL] menu={menu_no}, board_id={board_id}, attachments={len(attachments)}, "
            f"title={detail_meta.get('post_title', '')[:80]}"
        )

        if not attachments:
            return []

        records = []

        for index, attachment in enumerate(attachments, start=1):
            attachment_key = sha256_short(
                f"{menu_no}|{board_id}|{attachment.get('url', '')}|{attachment.get('name', '')}"
            )

            if attachment_key in existing_keys and not self.full:
                continue

            downloaded = self.download_attachment(attachment, post, index)

            if not downloaded:
                continue

            extracted_text = self.extract_text_from_file(downloaded["attachment_path"])
            extracted_text = clean_text(extracted_text)

            if extracted_text and self.is_garbled_text(extracted_text):
                print(f"[WARN] garbled text skipped: {downloaded['attachment_name']}")
                extracted_text = ""

            if not extracted_text:
                print(f"[WARN] no usable text extracted: {downloaded['attachment_name']}")

            final_text = clean_text(
                "\n\n".join(
                    [
                        f"[게시판 구분]\n{menu_name}",
                        f"[menuNo]\n{menu_no}",
                        f"[게시글 제목]\n{detail_meta.get('post_title', '')}",
                        f"[등록일]\n{detail_meta.get('registered_date', '')}",
                        f"[작성자]\n{detail_meta.get('author', '')}",
                        f"[게시글 URL]\n{detail_url}",
                        f"[첨부파일명]\n{downloaded['attachment_name']}",
                        "[게시글 본문 요약/안내문]",
                        post_body_text,
                        "[첨부파일 본문]",
                        extracted_text,
                    ]
                )
            )

            text_file_name = safe_filename(
                f"{menu_no}_{board_id}_{index:02d}_{detail_meta.get('post_title', '')}",
                220,
            )
            text_path = self.text_dir / safe_filename(f"{menu_no}_{menu_name}", 120) / f"{text_file_name}.txt"
            text_path.parent.mkdir(parents=True, exist_ok=True)
            text_path.write_text(final_text, encoding="utf-8")

            record = {
                "attachment_key": attachment_key,
                "menu_no": menu_no,
                "menu_name": menu_name,
                "post_uid": post_uid,
                "board_id": board_id,
                "post_title": detail_meta.get("post_title", ""),
                "registered_date": detail_meta.get("registered_date", ""),
                "author": detail_meta.get("author", ""),
                "detail_url": detail_url,
                "attachment_name": downloaded["attachment_name"],
                "attachment_url": downloaded["attachment_url"],
                "attachment_path": downloaded["attachment_path"],
                "attachment_text_path": str(text_path),
                "attachment_file_hash": downloaded["file_hash"],
                "attachment_text_hash": sha256_short(final_text),
                "attachment_text_chars": len(extracted_text),
                "post_body_text_chars": len(post_body_text),
                "rag_text_chars": len(final_text),
                "collected_at": datetime.now().isoformat(timespec="seconds"),
                "rag_text": final_text,
            }

            records.append(record)
            self.persist_record_incremental(record)

        return records

    def save_records(self, records: List[Dict]) -> None:
        records = sorted(
            records,
            key=lambda item: (
                item.get("registered_date", ""),
                item.get("menu_no", ""),
                item.get("board_id", ""),
                item.get("attachment_name", ""),
            ),
            reverse=True,
        )

        with self.records_jsonl.open("w", encoding="utf-8") as f:
            for record in records:
                f.write(json.dumps(record, ensure_ascii=False) + "\n")

        csv_fields = [
            "attachment_key",
            "menu_no",
            "menu_name",
            "post_uid",
            "board_id",
            "post_title",
            "registered_date",
            "author",
            "detail_url",
            "attachment_name",
            "attachment_url",
            "attachment_path",
            "attachment_text_path",
            "attachment_file_hash",
            "attachment_text_hash",
            "attachment_text_chars",
            "post_body_text_chars",
            "rag_text_chars",
            "collected_at",
        ]

        with self.records_csv.open("w", encoding="utf-8-sig", newline="") as f:
            writer = csv.DictWriter(f, fieldnames=csv_fields)
            writer.writeheader()

            for record in records:
                writer.writerow({key: record.get(key, "") for key in csv_fields})

        merged_parts = []

        for record in records:
            merged_parts.append(
                "\n".join(
                    [
                        "=" * 100,
                        f"menu_no: {record.get('menu_no', '')}",
                        f"menu_name: {record.get('menu_name', '')}",
                        f"board_id: {record.get('board_id', '')}",
                        f"post_title: {record.get('post_title', '')}",
                        f"registered_date: {record.get('registered_date', '')}",
                        f"attachment_name: {record.get('attachment_name', '')}",
                        f"detail_url: {record.get('detail_url', '')}",
                        "=" * 100,
                        record.get("rag_text", ""),
                    ]
                )
            )

        self.merged_txt.write_text("\n\n".join(merged_parts), encoding="utf-8")

        if self.collection_state is not None:
            try:
                self.collection_state.upsert_records(records)
                print(f"[SAVE] {self.state_db_path}")
            except Exception as exc:
                print(f"[WARN] failed to save collection state DB: {exc}")

        print(f"[SAVE] {self.records_jsonl}")
        print(f"[SAVE] {self.records_csv}")
        print(f"[SAVE] {self.merged_txt}")

    def run(self) -> None:
        if self.full and self.collection_state is not None:
            try:
                self.collection_state.clear()
                print(f"[FULL] cleared collection state DB: {self.state_db_path}")
            except Exception as exc:
                print(f"[WARN] failed to clear collection state DB: {exc}")

        existing_records = self.load_existing_records()
        existing_keys = self.load_existing_attachment_keys()

        posts = self.collect_all_posts()
        print(f"[INFO] total posts found: {len(posts)}")

        posts_with_attachment = [
            post for post in posts
            if post.get("has_attachment_icon", False)
        ]

        print(f"[INFO] total posts with attachment icon: {len(posts_with_attachment)}")

        new_records = []

        for post in tqdm(posts_with_attachment, desc="Collecting attachments"):
            records = self.collect_detail_attachments(post, existing_keys)
            new_records.extend(records)
            for record in records:
                key = record.get("attachment_key", "")
                if key:
                    existing_keys.add(key)

        all_records = self.load_existing_records() + existing_records + new_records

        dedup = {}
        for record in all_records:
            dedup[record["attachment_key"]] = record

        self.save_records(list(dedup.values()))

        print(f"[DONE] total attachment records: {len(dedup)}")
        print(f"[DONE] new attachment records: {len(new_records)}")

        if self.failed_jsonl.exists():
            print(f"[INFO] failed downloads saved to: {self.failed_jsonl}")


def main():
    parser = argparse.ArgumentParser()

    parser.add_argument(
        "--output-dir",
        type=str,
        default="data/sandan_attachment_kb",
    )
    parser.add_argument(
        "--sleep-sec",
        type=float,
        default=0.35,
    )
    parser.add_argument(
        "--timeout",
        type=int,
        default=40,
    )
    parser.add_argument(
        "--max-pages",
        type=int,
        default=300,
    )
    parser.add_argument(
        "--full",
        action="store_true",
    )
    parser.add_argument(
        "--show-pdf-warnings",
        action="store_true",
        help="Show pypdf internal warnings such as wrong pointing objects.",
    )
    parser.add_argument(
        "--no-pdf-repair",
        action="store_true",
        help="Disable optional PDF repair with pikepdf.",
    )
    parser.add_argument(
        "--state-db-path",
        type=str,
        default=None,
        help="SQLite path used to persist attachment collection state.",
    )

    args = parser.parse_args()

    collector = SandanAttachmentCollector(
        output_dir=args.output_dir,
        sleep_sec=args.sleep_sec,
        timeout=args.timeout,
        max_pages=args.max_pages,
        full=args.full,
        suppress_pdf_warnings=not args.show_pdf_warnings,
        repair_pdf=not args.no_pdf_repair,
        state_db_path=args.state_db_path,
    )

    collector.run()


if __name__ == "__main__":
    main()